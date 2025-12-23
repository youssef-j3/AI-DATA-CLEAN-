from pptx import Presentation
from pptx.util import Inches, Pt
import os

SRC = os.path.join(os.path.dirname(__file__), '..', 'slides', 'final_deck.md')
OUT = os.path.join(os.path.dirname(__file__), '..', 'slides', 'final_deck.pptx')


def text_to_slides(md_path, out_path):
    with open(md_path, 'r', encoding='utf-8') as f:
        content = f.read()
    # split slides by lines with only '---'
    parts = []
    current = []
    for line in content.splitlines():
        if line.strip() == '---':
            parts.append('\n'.join(current).strip())
            current = []
        else:
            current.append(line)
    if current:
        parts.append('\n'.join(current).strip())

    prs = Presentation()
    for i, slide_md in enumerate(parts):
        # use Title and Content layout when possible
        slide_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        lines = [ln for ln in slide_md.splitlines() if ln.strip()]
        if not lines:
            continue
        # first non-empty line as title
        title_text = lines[0]
        try:
            slide.shapes.title.text = title_text
        except Exception:
            pass
        body_lines = lines[1:]
        # try to find a content placeholder
        body = None
        for shape in slide.shapes:
            if not shape.is_placeholder:
                continue
            phf = shape.placeholder_format.type
            # 2 is body placeholder in many templates
            try:
                body = shape.text_frame
                break
            except Exception:
                body = None
        if body is None:
            # create a textbox
            left = Inches(1)
            top = Inches(1.5)
            width = Inches(8)
            height = Inches(4.5)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            body = txBox.text_frame
        for idx, bl in enumerate(body_lines):
            p = body.add_paragraph() if idx > 0 or body.text else body.paragraphs[0]
            p.text = bl
            p.level = 0
            p.font.size = Pt(18)
    prs.save(out_path)
    print(f"Saved PPTX to: {out_path}")


if __name__ == '__main__':
    text_to_slides(SRC, OUT)
